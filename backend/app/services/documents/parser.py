from __future__ import annotations

import csv
import io
import shutil
import subprocess
import tempfile
from collections.abc import Callable
from dataclasses import dataclass, field
from pathlib import Path

from app.core.config import get_settings


class DocumentParseError(RuntimeError):
    """A safe, user-facing document parsing failure."""


@dataclass(frozen=True)
class ParsedBlock:
    text: str
    page_number: int | None = None
    slide_number: int | None = None
    sheet_name: str | None = None
    heading_path: list[str] = field(default_factory=list)
    block_type: str = "paragraph"
    source_filename: str = ""
    char_start: int = 0
    char_end: int = 0


@dataclass(frozen=True)
class ParsedDocument:
    blocks: list[ParsedBlock]
    source_filename: str

    @property
    def text(self) -> str:
        return "\n\n".join(block.text for block in self.blocks if block.text.strip())

    # Compatibility for existing call sites while they migrate to structured blocks.
    def strip(self) -> str:
        return self.text.strip()

    def splitlines(self) -> list[str]:
        return self.text.splitlines()

    def split(self, *args, **kwargs) -> list[str]:
        return self.text.split(*args, **kwargs)

    def __len__(self) -> int:
        return len(self.text)

    def __str__(self) -> str:
        return self.text

    def __contains__(self, value: str) -> bool:
        return value in self.text


RawParser = Callable[[Path], list[dict]]
PARSERS: dict[str, RawParser] = {}


def register_parser(*suffixes: str):
    def decorator(parser: RawParser) -> RawParser:
        for suffix in suffixes:
            PARSERS[suffix.lower()] = parser
        return parser
    return decorator


def _block(text: str, **metadata) -> dict:
    return {"text": text.strip(), **metadata}


@register_parser(".pdf")
def parse_pdf(path: Path) -> list[dict]:
    try:
        import fitz
    except ImportError as exc:
        raise DocumentParseError("服务未安装 PyMuPDF，无法解析 PDF") from exc
    try:
        with fitz.open(path) as document:
            page_texts = [page.get_text("text").strip() for page in document]
            if sum(map(len, page_texts)) < get_settings().pdf_text_min_chars:
                if not get_settings().ocr_enabled:
                    raise DocumentParseError("PDF 未检测到可提取文本，且当前服务未启用 OCR")
                try:
                    import pytesseract
                    from PIL import Image
                    page_texts = []
                    for page in document:
                        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                        image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
                        page_texts.append(pytesseract.image_to_string(image, lang=get_settings().ocr_language).strip())
                except Exception as exc:
                    raise DocumentParseError(f"PDF OCR 失败，请确认 Tesseract 及语言包可用：{exc}") from exc
            return [_block(text, page_number=index, block_type="page") for index, text in enumerate(page_texts, 1) if text]
    except DocumentParseError:
        raise
    except Exception as exc:
        raise DocumentParseError(f"PDF 解析失败：{path.name}") from exc


@register_parser(".docx")
def parse_docx(path: Path) -> list[dict]:
    try:
        from docx import Document
        document = Document(path)
    except ImportError as exc:
        raise DocumentParseError("服务未安装 python-docx，无法解析 DOCX") from exc
    except Exception as exc:
        raise DocumentParseError(f"DOCX 解析失败：{path.name}") from exc
    blocks: list[dict] = []
    headings: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name if paragraph.style else "").lower()
        if style.startswith("heading"):
            try:
                level = max(1, int(style.split()[-1]))
            except ValueError:
                level = 1
            headings = headings[: level - 1] + [text]
            blocks.append(_block(text, heading_path=list(headings), block_type="heading"))
        else:
            blocks.append(_block(text, heading_path=list(headings), block_type="paragraph"))
    for table in document.tables:
        rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        text = "\n".join(row for row in rows if row.strip(" |"))
        if text:
            blocks.append(_block(text, heading_path=list(headings), block_type="table"))
    return blocks


@register_parser(".pptx")
def parse_pptx(path: Path) -> list[dict]:
    try:
        from pptx import Presentation
        presentation = Presentation(path)
    except ImportError as exc:
        raise DocumentParseError("服务未安装 python-pptx，无法解析 PPTX") from exc
    except Exception as exc:
        raise DocumentParseError(f"PPTX 解析失败：{path.name}") from exc
    blocks: list[dict] = []
    for index, slide in enumerate(presentation.slides, 1):
        seen: set[str] = set()
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text and text not in seen:
                    seen.add(text)
                    blocks.append(_block(text, slide_number=index, block_type="slide_text"))
            if getattr(shape, "has_table", False):
                text = "\n".join(" | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows)
                if text.strip() and text not in seen:
                    seen.add(text)
                    blocks.append(_block(text, slide_number=index, block_type="table"))
        notes = getattr(slide, "notes_slide", None)
        if notes is not None:
            note_text = "\n".join(shape.text.strip() for shape in notes.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip())
            if note_text and note_text not in seen:
                blocks.append(_block(note_text, slide_number=index, block_type="notes"))
    return blocks


@register_parser(".xlsx")
def parse_xlsx(path: Path) -> list[dict]:
    try:
        from openpyxl import load_workbook
        workbook = load_workbook(path, read_only=True, data_only=True)
    except ImportError as exc:
        raise DocumentParseError("服务未安装 openpyxl，无法解析 XLSX") from exc
    except Exception as exc:
        raise DocumentParseError(f"XLSX 解析失败：{path.name}") from exc
    try:
        blocks = []
        for sheet in workbook.worksheets:
            rows = [" | ".join(str(value) if value is not None else "" for value in row) for row in sheet.iter_rows(values_only=True)]
            text = "\n".join(row for row in rows if row.strip(" |"))
            if text:
                blocks.append(_block(text, sheet_name=sheet.title, block_type="sheet"))
        return blocks
    finally:
        workbook.close()


def _decode_bytes(data: bytes, path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise DocumentParseError(f"文本编码无法识别：{path.name}")


@register_parser(".csv")
def parse_csv(path: Path) -> list[dict]:
    text = _decode_bytes(path.read_bytes(), path)
    return [_block("\n".join(" | ".join(row) for row in csv.reader(io.StringIO(text))), block_type="table")]


@register_parser(".txt", ".md")
def parse_text(path: Path) -> list[dict]:
    return [_block(_decode_bytes(path.read_bytes(), path), block_type="text")]


def _parse_legacy_office(path: Path) -> list[dict]:
    target_suffix = {".doc": ".docx", ".ppt": ".pptx", ".xls": ".xlsx"}[path.suffix.lower()]
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise DocumentParseError("当前服务未安装 LibreOffice/soffice，无法解析旧版 Office 文件")
    with tempfile.TemporaryDirectory(prefix="prismmind-office-") as temp_dir:
        try:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", target_suffix[1:], "--outdir", temp_dir, str(path)],
                capture_output=True, text=True, timeout=get_settings().office_conversion_timeout_seconds, check=False,
            )
        except subprocess.TimeoutExpired as exc:
            raise DocumentParseError("旧版 Office 文件转换超时") from exc
        converted = Path(temp_dir) / f"{path.stem}{target_suffix}"
        if result.returncode != 0 or not converted.exists():
            detail = (result.stderr or result.stdout or "未知错误").strip()
            raise DocumentParseError(f"LibreOffice 转换失败：{detail[:300]}")
        return PARSERS[target_suffix](converted)


for _legacy_suffix in (".doc", ".ppt", ".xls"):
    PARSERS[_legacy_suffix] = _parse_legacy_office


def parse_document(file_path: str | Path, suffix: str | None = None) -> ParsedDocument:
    path = Path(file_path)
    normalized_suffix = (suffix or path.suffix).lower()
    parser = PARSERS.get(normalized_suffix)
    if parser is None:
        raise DocumentParseError(f"不支持的文档格式：{normalized_suffix or '无后缀'}")
    raw_blocks = parser(path)
    blocks: list[ParsedBlock] = []
    cursor = 0
    for raw in raw_blocks:
        text = str(raw.pop("text", "")).strip()
        if not text:
            continue
        start = cursor
        end = start + len(text)
        blocks.append(ParsedBlock(text=text, source_filename=path.name, char_start=start, char_end=end, **raw))
        cursor = end + 2
    if not blocks:
        raise DocumentParseError("文档解析完成，但未提取到任何文本")
    return ParsedDocument(blocks=blocks, source_filename=path.name)


def parser_name_for_suffix(suffix: str) -> str:
    parser = PARSERS.get(suffix.lower())
    return parser.__name__ if parser else "unsupported"
